# Import PyMuPDF for opening PDF files and extracting text
import pymupdf

# Import re for cleaning and processing extracted text
import re

# Import Path for working with file paths
from pathlib import Path

# Import python-docx for opening DOCX files and extracting text
from docx import Document

# Import python-pptx for opening PPTX files and extracting text
from pptx import Presentation

# Function to clean extracted text
def clean_text(text):

    # Remove excessive spaces
    text = re.sub(r"[ \t]+", " ", text)         # Replace multiple spaces and tabs with a single space

    # Remove excessive blank lines
    text = re.sub(r"\n\s*\n+", "\n\n", text)   # Replace multiple blank lines with a single blank line

    # Fix spaces before punctuation
    text = re.sub(r"\s+([,.!?;:])", r"\1", text)   # Remove spaces before punctuation marks

    # Remove spaces from the beginning and end
    return text.strip()                         # Remove leading and trailing spaces from the text




#-------------------------------------------------------------PDF Processing Functions-------------------------------------------------------------


# Function to extract text from a PDF file
def extract_text_from_the_pdf(pdf_path):                # Function to extract text from a PDF file

    # Open the PDF file
    doc = pymupdf.open(pdf_path)                   # Open the PDF file using PyMuPDF

    # Create an empty list to store text from each page
    pages = []

    # Loop through every page in the PDF
    for page_number, page in enumerate(doc, start=1):   # Enumerate through each page in the PDF starting from page 1

        # Extract text from the current page
        text = page.get_text()                         # Extract text from the current page using PyMuPDF's get_text() method

        # Clean the extracted text
        text = clean_text(text)                        # Clean the extracted text using the clean_text function

        # Store the page number and extracted text
        pages.append({
            "page": page_number,   # Page Number
            "text": text            # Text
        })

    # Close the PDF file
    doc.close()             # Close File After Extraction

    # Return the extracted text
    return pages


#-------------------------------------------------------------DOCX Processing Functions-------------------------------------------------------------


# Function to extract text from a DOCX file
def extract_text_from_docx(file_path):                # Function to extract text from a DOCX file

    # Open the DOCX file
    doc = Document(file_path)                         # Open the DOCX file using python-docx

    # Create an empty list to store text from each paragraph
    paragraphs = []

    # Loop through every paragraph in the DOCX file
    for paragraph_number, paragraph in enumerate(doc.paragraphs, start=1):   # Enumerate through each paragraph starting from paragraph 1

        # Extract text from the current paragraph
        text = paragraph.text                         # Extract text from the current paragraph

        # Clean the extracted text
        text = clean_text(text)                       # Clean the extracted text using the clean_text function

        # Skip empty paragraphs
        if not text:
            continue

        # Store the paragraph number and extracted text
        paragraphs.append({
            "paragraph": paragraph_number,   # Paragraph Number
            "text": text                     # Text
        })

    # Return the extracted paragraphs
    return paragraphs



#-------------------------------------------------------------TXT Processing Functions-------------------------------------------------------------


# Function to extract text from a TXT file
def extract_text_from_txt(file_path):                # Function to extract text from a TXT file

    # Open the TXT file
    with open(file_path, "r", encoding="utf-8", errors="ignore") as txt_file:   # Open the TXT file and ignore encoding issues

        # Read the text from the TXT file
        text = txt_file.read()                        # Read all the text from the TXT file

    # Clean the extracted text
    text = clean_text(text)                           # Clean the extracted text using the clean_text function

    # Create an empty list to store text paragraphs
    paragraphs = []

    # Loop through every paragraph in the text
    for paragraph_number, paragraph in enumerate(re.split(r"\n\s*\n+", text), start=1):   # Split the text by blank lines

        # Clean the paragraph text
        paragraph = clean_text(paragraph)              # Clean the paragraph text

        # Skip the paragraph if no text is available
        if not paragraph:
            continue

        # Store the paragraph number and extracted text
        paragraphs.append({
            "paragraph": paragraph_number,   # Paragraph Number
            "text": paragraph                # Paragraph Text
        })

    # Return the extracted paragraphs
    return paragraphs


##-------------------------------------------------------------text processing from ppt-------------------------------------------------------------



# Function to extract text from a PPTX file
def extract_text_from_pptx(file_path):       # Function to extract text from a PPTX file

    # Open the PPTX file
    presentation = Presentation(file_path)   # Open the PPTX file using python-pptx

    # Create an empty list to store text from each slide
    slides = []

    # Loop through every slide in the presentation
    for slide_number, slide in enumerate(presentation.slides, start=1):

        # Create an empty string to store slide text
        slide_text = ""

        # Loop through every shape on the slide
        for shape in slide.shapes:

            # Check if the shape contains text
            if hasattr(shape, "text"):

                # Add the shape text to the slide text
                slide_text += shape.text + "\n"

        # Clean the extracted text
        slide_text = clean_text(slide_text)   # Clean the extracted slide text

        # Skip the slide if no text is available
        if not slide_text:
            continue

        # Store the slide number and extracted text
        slides.append({
            "slide": slide_number,   # Slide Number
            "text": slide_text       # Slide Text
        })

    # Return the extracted slides
    return slides

#-------------------------------------------------------------Chunking Functions-------------------------------------------------------------


# Function to create chunks of text
def create_chunks(text, chunk_size=1000, overlap=200):   # Function to create chunks of text with specified size and overlap

    # Split the text into individual words
    words = text.split()   # Split the text into a list of individual words

    # Create an empty list to store the chunks of text
    chunks = []             # Empty list to store the chunks of text

    # Initialize the starting index for chunking
    start = 0               # Initialize the starting index for chunking

    # Loop until all words are processed
    while start < len(words):   # Loop until all words are processed

        # Calculate the ending index for the current chunk
        end = start + chunk_size   # Calculate the ending index for the current chunk

        # Join the words from the starting index to the ending index
        chunk = " ".join(words[start:end])   # Join the selected words to form a chunk of text

        # Add the current chunk to the list of chunks
        chunks.append(chunk)   # Add the current chunk to the list of chunks

        # Move the starting index forward while maintaining overlap
        start += chunk_size - overlap   # Move forward by the chunk size minus the overlap to maintain context

    # Return the created chunks
    return chunks


#-------------------------------------------------------------Processed Chunks Creation Functions-------------------------------------------------------------



# Function to create processed chunks from extracted text
def create_processed_chunks(items, source):       # Function to create processed chunks from extracted text

    # Create an empty list to store processed chunks
    processed_chunks = []

    # Loop through every extracted item
    for item in items:

        # Get the extracted text from the current item
        text = item["text"]

        # Skip the item if no text is available
        if not text:
            continue

        # Create chunks from the extracted text
        chunks = create_chunks(text)

        # Loop through every chunk
        for chunk_number, chunk in enumerate(chunks, start=1):

            # Store the source, chunk number, and text
            processed_chunk = {
                "source": source,        # Source file name
                "chunk": chunk_number,   # Chunk Number
                "text": chunk             # Chunk Text
            }

            # Add the page number if available
            if "page" in item:

                processed_chunk["page"] = item["page"]   # Page Number

            # Add the paragraph number if available
            if "paragraph" in item:

                processed_chunk["paragraph"] = item["paragraph"]   # Paragraph Number

            # Add the processed chunk to the list
            processed_chunks.append(processed_chunk)

    # Return all processed chunks
    return processed_chunks



#-------------------------------------------------------------File Processing Functions (PDF)-------------------------------------------------------------



# Function to process the complete PDF
def process_pdf(pdf_path):       # Function to process a PDF

    # Extract text from the PDF
    pages = extract_text_from_the_pdf(pdf_path)

    # Get the PDF file name
    source = Path(pdf_path).name

    # Create processed chunks from the extracted pages
    processed_chunks = create_processed_chunks(pages, source)

    # Return all processed chunks
    return processed_chunks


#-------------------------------------------------------------DOCX Processing Functions-------------------------------------------------------------


# Function to process the complete DOCX file
def process_docx(file_path):       # Function to process a DOCX file

    # Extract text from the DOCX file
    paragraphs = extract_text_from_docx(file_path)

    # Get the DOCX file name
    source = Path(file_path).name

    # Create processed chunks from the extracted paragraphs
    processed_chunks = create_processed_chunks(paragraphs, source)

    # Return all processed chunks
    return processed_chunks



#-------------------------------------------------------------TXT Processing Functions-------------------------------------------------------------


# Function to process the complete TXT file
def process_txt(file_path):       # Function to process a TXT file

    # Extract text from the TXT file
    paragraphs = extract_text_from_txt(file_path)

    # Get the TXT file name
    source = Path(file_path).name

    # Create processed chunks from the extracted paragraphs
    processed_chunks = create_processed_chunks(paragraphs, source)

    # Return all processed chunks
    return processed_chunks


#-----------------------------------------------------PPTX Processing Functions

# Function to process the complete PPTX file
def process_pptx(file_path):       # Function to process a PPTX file

    # Extract text from the PPTX file
    slides = extract_text_from_pptx(file_path)

    # Get the PPTX file name
    source = Path(file_path).name

    # Create processed chunks from the extracted slides
    processed_chunks = create_processed_chunks(slides, source)

    # Return all processed chunks
    return processed_chunks

#-------------------------------------------------------------File Processing Functions (Automatic) Detection------------------------------------------------------------


# Function to process different file formats
def process_file(file_path):       # Function to process a file according to its format

    # Get the file extension
    file_extension = Path(file_path).suffix.lower()   # Get the file extension and convert it to lowercase

    # Check if the file is a PDF
    if file_extension == ".pdf":

        # Process the PDF file
        return process_pdf(file_path)   # Send the PDF file to the PDF processing function

    # Check if the file is a DOCX file
    elif file_extension == ".docx":

        # Process the DOCX file
        return process_docx(file_path)   # Send the DOCX file to the DOCX processing function

    # Check if the file is a TXT file
    elif file_extension == ".txt":

        # Process the TXT file
        return process_txt(file_path)   # Send the TXT file to the TXT processing function

    # Check if the file is a PPTX file
    elif file_extension == ".pptx":

        # Process the PPTX file
        return process_pptx(file_path)   # Send the PPTX file to the PPTX processing function

    # If the file format is not supported
    else:

        # Display an error message
        raise ValueError("Unsupported file format. Only PDF, DOCX, TXT , And PPTX files are supported.")



#-------------------------------------------------------------Example Usage-------------------------------------------------------------



# Temporary File Path
file_path = "---------"  # Temporary File Path


# Process the file automatically according to its format
processed_chunks = process_file(file_path)   # Automatically select the correct processing function


# Print the number of processed chunks
print("Number of processed chunks:", len(processed_chunks))   # This prints the number of processed chunks


# Check if processed chunks are available
if processed_chunks:

    # Print the first processed chunk
    print("\nFirst processed chunk:")   # Print the first processed chunk

    print(processed_chunks[0])   # Display the first processed chunk

    # Print the second processed chunk if available
    if len(processed_chunks) > 1:

        print("\nSecond processed chunk:")   # Print the second processed chunk

        print(processed_chunks[1])   # Display the second processed chunk

    # Print the first three processed chunks
    print("\nFirst three processed chunks:")

    for item in processed_chunks[:3]:

        print("\n------------------------")

        print("Source:", item["source"])   # Print the source of the chunk

        # Check whether the chunk belongs to a PDF
        if "page" in item:

            print("Page:", item["page"])   # Print the page number

        # Check whether the chunk belongs to a DOCX or TXT file
        if "paragraph" in item:

            print("Paragraph:", item["paragraph"])   # Print the paragraph number

        print("Chunk:", item["chunk"])   # Print the chunk number

        print(item["text"][:500])   # Print the first 500 characters of the chunk


def retrieve_chunks(question, processed_chunks, top_k=5):

    question_words = set(question.lower().split())

    scored_chunks = []

    for chunk in processed_chunks:

        text_words = set(chunk["text"].lower().split())

        score = len(question_words.intersection(text_words))

        scored_chunks.append((score, chunk))

    scored_chunks.sort(key=lambda x: x[0], reverse=True)

    retrieved_chunks = []

    for score, chunk in scored_chunks[:top_k]:

        if score > 0:
            retrieved_chunks.append(chunk)

    return retrieved_chunks
        