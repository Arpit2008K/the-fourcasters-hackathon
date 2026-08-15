import streamlit as st      # For streamlit web app interface
import pymupdf              # For PDF text extraction
from docx import Document   # For DOCX text extraction
from pptx import Presentation   # For PPTX text extraction
import io                       # For handling in-memory file operations

# Text extraction function for various file types
def extract_text(uploaded_file):
    # Get the file name and extension from the uploaded file
    file_name = uploaded_file.name
    extension = file_name.split(".")[-1].lower()

    try:
        # For PDF files, use PyMuPDF to extract text
        if extension == "pdf":
            doc = pymupdf.open(
                stream=uploaded_file.getvalue(),        # Open the PDF file from the uploaded file's byte stream
                filetype="pdf"
            )
            text = ""
            for page_number, page in enumerate(doc, start=1):   # Loop through each page in the PDF starting from page 1
                page_text = page.get_text()                     # Extract the text from the current page

                text += f"\n\nPage {page_number}\n"             # Add page number header to the extracted text
                text += page_text
            doc.close()                                 # Close the PDF file after extraction
            return text


        # For DOCX files, use python-docx to extract text
        elif extension == "docx":

            document = Document(
                io.BytesIO(uploaded_file.getvalue())        # Open the DOCX file from the uploaded file's byte stream
            )
            text = ""
            for paragraph in document.paragraphs:       # To extract the text and append it to the text variable
                text += paragraph.text + "\n"
            return text


        # For PPTX files, use python-pptx to extract text
        elif extension == "pptx":
            presentation = Presentation(
                io.BytesIO(uploaded_file.getvalue())        # Open the PPTX file from the uploaded file's byte stream
            )
            text = ""
            # Loop through each slide in the presentation starting from slide 1
            for slide_number, slide in enumerate(presentation.slides,start=1):          
                text += f"\n\nSlide {slide_number}\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):          # Check if the shape has text attribute
                        text += shape.text + "\n"
            return text

        # For txt files, decode the bytes to a string
        elif extension == "txt":
            return uploaded_file.getvalue().decode(
                "utf-8",                    # Decode the byte stream to a UTF-8 string
                errors="ignore"
            )
        else:
            return "Unsupported file type."
        
    except Exception as e:

        return f"ERROR: {str(e)}"       # Return an error message if any exception occurs during text extraction

# Streamlit app configuration 
st.set_page_config(
    page_title="AI Research Assistant",
    page_icon="📚"
)

st.title(" _AI Research Assistant_")

uploaded_files = st.file_uploader(
    "Upload your research files",
    type=[
        "pdf",
        "docx",
        "pptx",
        "txt",
    ],
    accept_multiple_files=True)

if uploaded_files:

    st.success(
        f"{len(uploaded_files)} file(s) uploaded successfully!"
    )

    for uploaded_file in uploaded_files:
        st.subheader(f"{uploaded_file.name}")           # Display the name of each uploaded file as a subheader
        extracted_text = extract_text(uploaded_file)

        if extracted_text.startswith("ERROR:"):
            st.error(extracted_text)        # Error message if text extraction fails
        elif extracted_text.strip():        # Check if the extracted text is not empty
            st.success(
                f"Text extracted successfully — "
                f"{len(extracted_text)} characters"
            )

            with st.expander("View extracted text"):        # to view the extracted text in an expandable section
                st.text_area(
                    "Extracted Content",
                    extracted_text,
                    height=400,
                    key=f"text_{uploaded_file.name}"
                )

        else:
            st.warning(
                "No text could be extracted from this file."
            )